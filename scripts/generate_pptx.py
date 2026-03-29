import json
import sys
import os
import uuid
from io import BytesIO
from urllib.request import urlopen, Request

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw


class PptxPresentationCreator:
    def __init__(self, ppt_json_model):
        self.ppt_model = ppt_json_model
        self.ppt = Presentation()
        # Force a 16:9 slide size (approximate using points)
        self.ppt.slide_width = Pt(1280)
        self.ppt.slide_height = Pt(720)

    def generate(self, output_path: str):
        for slide_data in self.ppt_model.get("slides", []):
            slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[6])  # blank

            bg_color = slide_data.get("background_color")
            if bg_color:
                # Expect hex string without '#'
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor.from_string(bg_color)

            for shape in slide_data.get("shapes", []):
                shape_type = shape.get("type")
                if shape_type == "text":
                    self.add_textbox(slide, shape)
                elif shape_type == "image":
                    self.add_picture(slide, shape)

        self.ppt.save(output_path)

    def add_textbox(self, slide, text_shape):
        position = text_shape.get("position", {})
        left = Pt(position.get("left", 0))
        top = Pt(position.get("top", 0))
        width = Pt(position.get("width", 400))
        height = Pt(position.get("height", 100))

        textbox_shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox_shape.text_frame
        text_frame.clear()

        paragraphs = text_shape.get("paragraphs", [])
        first = True
        for paragraph_data in paragraphs:
            if first:
                p = text_frame.paragraphs[0]
                first = False
            else:
                p = text_frame.add_paragraph()
            for run_data in paragraph_data.get("runs", []):
                run = p.add_run()
                run.text = run_data.get("text", "")
                font_size = run_data.get("font_size")
                if font_size is not None:
                    run.font.size = Pt(font_size)
                font_name = run_data.get("font_name")
                if font_name:
                    run.font.name = font_name
                color_hex = run_data.get("color_hex")
                if color_hex:
                    run.font.color.rgb = RGBColor.from_string(color_hex)
                if run_data.get("bold"):
                    run.font.bold = True
                if run_data.get("italic"):
                    run.font.italic = True

            # If this textbox is flagged as bullets, ensure bullet-style spacing
            if text_shape.get("is_bullets"):
                p.level = 0
                try:
                    # Slight extra spacing between bullet lines
                    p.space_after = Pt(4)
                except Exception:
                    pass

    def add_picture(self, slide, image_shape):
        position = image_shape.get("position", {})
        left = Pt(position.get("left", 0))
        top = Pt(position.get("top", 0))
        width = Pt(position.get("width", 400))
        height = Pt(position.get("height", 300))

        url = image_shape.get("url")
        if not url:
            return

        try:
            req = Request(url, headers={"User-Agent": "Mozilla/5.0"})
            with urlopen(req) as resp:
                data = resp.read()
        except Exception as e:
            # If image download fails, skip this image but keep the rest of the slide.
            sys.stderr.write(f"Image download failed for {url}: {e}\n")
            return

        im = Image.open(BytesIO(data)).convert("RGBA")

        # Apply rounded corners if specified in the JSON model
        radius = image_shape.get("border_radius")
        if radius:
            im = self.round_image_corners(im, int(radius))
        tmp_dir = os.path.join("/tmp", "presenton_pptx_imgs")
        os.makedirs(tmp_dir, exist_ok=True)
        tmp_img_path = os.path.join(tmp_dir, f"{uuid.uuid4().hex}.png")
        im.save(tmp_img_path)

        slide.shapes.add_picture(tmp_img_path, left, top, width=width, height=height)

    def round_image_corners(self, im, rad):
        """Return a copy of the image with rounded corners (alpha mask)."""
        w, h = im.size
        radius = max(0, min(rad, min(w, h) // 2))
        if radius == 0:
            return im

        mask = Image.new("L", (w, h), 0)
        draw = ImageDraw.Draw(mask)
        draw.rounded_rectangle((0, 0, w, h), radius=radius, fill=255)
        im = im.copy()
        im.putalpha(mask)
        return im


def main():
    if len(sys.argv) != 3:
        print("Usage: generate_pptx.py <model_json_path> <output_pptx_path>", file=sys.stderr)
        sys.exit(1)

    model_path = sys.argv[1]
    output_path = sys.argv[2]

    with open(model_path, "r", encoding="utf-8") as f:
        model = json.load(f)

    creator = PptxPresentationCreator(model)
    creator.generate(output_path)


if __name__ == "__main__":
    main()
